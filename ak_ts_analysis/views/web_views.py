import os
import json
import tempfile
from pptx import Presentation
from django.http import JsonResponse, HttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.shortcuts import render
from django.conf import settings
import zipfile
import xml.etree.ElementTree as ET

import olefile


def extract_excel_from_ole(bin_path, output_dir):
    """
    Попытка извлечь Excel из OLE контейнера (.bin).
    Если удачно — сохраняет .xlsx рядом и возвращает путь к нему,
    иначе возвращает None.
    """
    if not olefile.isOleFile(bin_path):
        return None

    ole = olefile.OleFileIO(bin_path)

    # Часто Excel-объект лежит в потоке "Workbook" или "Package"
    possible_streams = ['Workbook', 'Package']

    for stream_name in possible_streams:
        if ole.exists(stream_name):
            data = ole.openstream(stream_name).read()
            # Сохраняем как .xlsx
            output_path = os.path.splitext(bin_path)[0] + '.xlsx'
            with open(output_path, 'wb') as f:
                f.write(data)
            ole.close()
            return output_path

    ole.close()
    return None


def extract_embedded_excels_with_slide_map(pptx_path, output_dir):
    """
    Извлекает встроенные Excel-файлы из .pptx, сохраняет в output_dir,
    возвращает dict { slide_number: [список excel файлов для этого слайда] }
    При этом из .bin файлов извлекает .xlsx с помощью extract_excel_from_ole.
    """
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    slide_excel_map = {}  # slide_number -> list of excel filenames
    extracted_files = set()
    extracted_bin_excel_files = {}

    with zipfile.ZipFile(pptx_path, 'r') as zipf:
        # Получаем список слайдов (например: ppt/slides/slide1.xml, slide2.xml и т.д.)
        slide_files = sorted([f for f in zipf.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')])

        for slide_file in slide_files:
            # Определяем номер слайда из имени файла
            slide_num = int(''.join(filter(str.isdigit, slide_file)))

            # Файл rels для этого слайда
            rels_path = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
            excel_files_for_slide = []

            if rels_path in zipf.namelist():
                rels_data = zipf.read(rels_path)
                root = ET.fromstring(rels_data)

                # Пространство имён для rels
                ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}

                for rel in root.findall('r:Relationship', ns):
                    target = rel.attrib.get('Target', '')
                    if target.startswith('../embeddings/') and (
                        target.endswith('.xlsx') or target.endswith('.xls') or target.endswith('.bin')
                    ):
                        # Формируем путь внутри zip
                        embedded_path = 'ppt/' + target.replace('../', '')
                        filename = os.path.basename(embedded_path)

                        # Если еще не извлекали — извлекаем
                        if filename not in extracted_files:
                            extracted_files.add(filename)
                            output_path = os.path.join(output_dir, filename)
                            with open(output_path, 'wb') as out_file:
                                out_file.write(zipf.read(embedded_path))

                            # Если .bin — пробуем извлечь из OLE Excel
                            if filename.endswith('.bin'):
                                excel_path = extract_excel_from_ole(output_path, output_dir)

                                if excel_path:
                                    print(excel_path)
                                    excel_filename = os.path.basename(excel_path)
                                    extracted_files.add(excel_filename)
                                    extracted_bin_excel_files[filename] = excel_filename
                                    excel_files_for_slide.append(excel_filename)
                                else:
                                    pass
                                    # Не удалось извлечь excel, добавим .bin файл как есть
                                    #excel_files_for_slide.append(filename)
                            else:
                                excel_files_for_slide.append(filename)

                            try:
                                os.remove(output_path)
                            except Exception as e:
                                print(f'remove error {e}')
                        else:
                            # Уже извлекали, просто добавляем имя
                            if '.xls' in filename:
                                excel_files_for_slide.append(filename)
                            elif filename in extracted_bin_excel_files:
                                excel_files_for_slide.append(filename)

            slide_excel_map[slide_num] = excel_files_for_slide

    return slide_excel_map


@csrf_exempt
def upload_presentation(request):
    if request.method == 'POST':
        ppt_file = request.FILES.get('presentation')
        if not ppt_file:
            return HttpResponse('Файл презентации не был загружен', status=400)

        # Сохраняем файл во временное хранилище
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            for chunk in ppt_file.chunks():
                tmp.write(chunk)
            pptx_path = tmp.name

        # Создаём рабочую директорию для вывода
        base_name = os.path.splitext(ppt_file.name)[0]
        output_dir = os.path.join(settings.MEDIA_ROOT, 'pptx_outputs', f"output_{base_name}")
        os.makedirs(output_dir, exist_ok=True)

        prs = Presentation(pptx_path)

        # Извлекаем Excel-файлы с привязкой к слайдам, включая обработку .bin -> .xlsx
        slide_excel_map = extract_embedded_excels_with_slide_map(pptx_path, output_dir)

        slide_data = []

        for idx, slide in enumerate(prs.slides):
            slide_index = idx + 1
            title_text = ""
            tables = []

            # Получение заголовка
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text.strip()

            # Обход всех объектов слайда
            for shape in slide.shapes:
                try:
                    if shape.shape_type == 19:  # таблица
                        table = []
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table.append(row_data)
                        tables.append(table)
                except NotImplementedError:
                    continue

            embedded_excels = slide_excel_map.get(slide_index, [])

            slide_data.append({
                "slide": slide_index,
                "title": title_text,
                "tables": tables,
                "embedded_excels": embedded_excels
            })

        # Сохраняем JSON
        json_path = os.path.join(output_dir, f"output_{base_name}.json")
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(slide_data, f, ensure_ascii=False, indent=2)

        # Удаляем временный pptx файл
        os.remove(pptx_path)

        return JsonResponse({
            "base_name": base_name,
            "slides": slide_data
        }, json_dumps_params={"ensure_ascii": False, "indent": 2})

    return render(request, 'ak_ts_analysis/upload.html')
